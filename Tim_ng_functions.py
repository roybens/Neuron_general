import itertools
import os
import json
from PyPDF2 import PdfMerger #need to write for pypdf
from pypdf import PdfWriter
from pypdf import PdfReader
import pandas as pd
import csv
import re
from pptx import Presentation
from pptx.util import Inches,Pt
from PIL import Image
import io
import fitz
import sys
import numpy as np
import matplotlib.pyplot as plt
from Na12HMMModel_TF import *

def combine_pdfs(folder_path, out_sfx): #input folder pather where pdfs are stored, out_sfx = output suffix
    file = open('./JUPYTERmutant_list.txt','r') #put mutant names in here
    Lines = file.readlines() #read each mut name
    #os.chdir(folder_path)
    for line in Lines:
        mutTXT = line.strip() #remove return after each mut name
        print(mutTXT)
        mutfolderpath = os.path.join(folder_path,mutTXT)
        
        os.chdir(mutfolderpath)
        #
        #if filename.endswith(".pdf"):
            #merger.append(os.path.join(mutfolderpath, filename))
        merger = PdfMerger()
        #merger = PdfWriter() #for use with pypdf library *UNTESTED*

        # pdfs = [f'{mutTXT}__0.5_75.pdf',f'{mutTXT}__0.5_100.pdf',
        #         'test__fi.pdf','test_0.5_wtvmut.pdf','test_1_wtvmut.pdf']
        pdfs = ['test_0.5_wtvmut.pdf','test_1_wtvmut.pdf']
        
        for pdf in pdfs:
            merger.append(pdf)
        
        #merger.write(folder_path+"_combined100323.pdf")
        merger.write(folder_path+mutTXT+out_sfx+'.pdf')
    merger.close()

# def convert_pdf_to_jpg(pdf_path, output_jpg_path):
#     """
#     Converts a single page of a PDF to a JPG image.

#     Args:
#         pdf_path: Path to the PDF file.
#         output_jpg_path: Path to the output JPG image.
#     """

#     # Open the PDF and get the first page
#     reader = PdfReader(pdf_path)
#     page = reader.getPage(0)

#     # Extract the page content as a PIL Image object
#     pil_image = Image.open(page.extractImage()[0])

#     # Convert the image to RGB mode (required for JPG)
#     pil_image = pil_image.convert("RGB")

#     # Save the image as a JPG
#     pil_image.save(output_jpg_path, quality=95)

#     # Example usage
#     pdf_path = "/path/to/your.pdf"
#     output_jpg_path = "/path/to/output.jpg"
# convert_pdf_to_jpg(pdf_path, output_jpg_path)

# print(f"Successfully converted PDF page to JPG: {output_jpg_path}")


def combine_and_sort_ef_csvs(root_path, out_sfx):
  combined_ef = f'{root_path}combined_ef_{out_sfx}.csv'
  # Open the output file in write mode
  with open(combined_ef, "w", newline="") as outfile:
    writer = csv.writer(outfile)

    data_by_filename = {}

    # Loop through all folders in the current directory
    for folder in os.listdir(root_path):
      # Check if it's a directory and not a file
      if os.path.isdir(os.path.join(root_path, folder)):
        # Access files within the folder using another loop
        for filename in os.listdir(os.path.join(root_path, folder)):
          if filename.endswith(".csv"):
            # Extract mutation number and index from filename
            match = re.search(r"mut(\d+)_(\d+)_", filename)
            if match:
              mutation_number, index = match.groups()

              # Open and process the CSV file
              with open(os.path.join(root_path, folder, filename), "r", newline="") as infile:
                reader = csv.reader(infile)

                # Skip the header
                next(reader)

                # Store data row with filename information as key
                for row in reader:
                  data_by_filename[f"{mutation_number}_{index}_{filename}"] = row

    # Sort data based on filename key (automatically sorts based on structure)
    sorted_data = sorted(data_by_filename.items(), key=lambda item: item[0])

    # Write sorted data to the output file
    for key, row in sorted_data:
      writer.writerow(row)

  print(f"Successfully combined and sorted data into '{combined_ef}'!")
  return


def make_ppt_from_pdf2(pdf_path, output_ppt_path):
  prs = Presentation()
  blank_slide_layout = prs.slide_layouts[6]
  slide = prs.slides.add_slide(blank_slide_layout)
  left = top = width = height = Inches(0.5)
  
  
  # file = open("/global/homes/t/tfenton/Neuron_general-2/JUPYTERmutant_list.txt", "r")
  # lines = file.readlines()
  # for line in lines: #mut1_1,mut1_2 etc.
  #   mut_txt = line.strip()
  #   print(mut_txt)
    
  for folder in sorted(os.listdir(pdf_path)): #Sort the listed directories so output pptx is somewhat in order 
      if os.path.isdir(os.path.join(pdf_path, folder)):# Check if it's a directory
        slide = prs.slides.add_slide(blank_slide_layout) #Add slide for each folder. Each slide will have multiple plots on it
        txBox = slide.shapes.add_textbox(left,top,width,height) #Add text box for title
        tf = txBox.text_frame
        page_count = 0
        for filename in sorted(os.listdir(os.path.join(pdf_path, folder))):# Access files within the folder, and sort them to appear in sam order on slide
          if filename.endswith(".pdf"):
  
            file = os.path.join(pdf_path,folder, filename) #pdf file path
            
            tf.text = folder
            print(folder)
            
            doc = fitz.open(file)
            for page in doc:  # iterate through the pages
              pix = page.get_pixmap(dpi=150)  # render page to an image
              pix.save(f"{pdf_path}/{folder}-{filename}.png")  # store image as a PNG
              
              
              slide.shapes.add_picture(f"{pdf_path}/{folder}-{filename}.png",left=Inches(page_count*2.5),top=Inches(1), width=Inches(2.5))
              print(page_count)
          page_count+=1
          
  
  prs.save(output_ppt_path)
  print(f"Successfully converted PDFs to PowerPoint presentation: {output_ppt_path}")
  return


#This function plots efel efeatures as bar graphs (in this case HH and HMM)
def plot_efeatures_bar(plot_folder,pfx):
    x = ['HH','HMM']
    data1 = np.genfromtxt('na12_orig1_efel.csv',delimiter=',')
    data2 = np.genfromtxt('na12_HMM_TF100923_efel.csv',delimiter=',')
    for i in range(1,15):
      y = [data1[1,i],data2[1,i]]
      print(y)
      fig, ax = plt.subplots()
      ax.bar(x,y,width=0.3,edgecolor='white')
      file_path_to_save=f'{plot_folder}{pfx}_{i}.pdf'
      plt.savefig(file_path_to_save, format='pdf')

    return


##Takes params text file and allows you to change the values for scanning etc...
def modify_dict_file(filename, changes):
  """
  Modifies values in a dictionary stored in a text file.

  Args:
      filename: The name of the text file containing the dictionary.
      changes: A dictionary containing key-value pairs where the key is the key to modify in the original dictionary and the value is the new value.

  Raises:
      ValueError: If the file cannot be opened or the content is not valid JSON.
  """

  try:
    # Open the file and read its content
    with open(filename, "r") as file:
      content = file.read()

    # Try to load the content as a dictionary
    try:
      data = eval(content)  # Assuming the file contains valid dictionary syntax
    except (NameError, SyntaxError):
      raise ValueError("Invalid dictionary format in the file.")

    # Modify values based on the provided changes dictionary
    for key, value in changes.items():
      if key not in data:
        print(f"Warning: Key '{key}' not found in the dictionary, skipping.")
      else:
        data[key] = value

    # Write the modified dictionary back to the file
    # with open(filename, "w") as file:
    #   file.write(repr(data))
    with open(filename, "w") as file:
      file.write(json.dumps(data, indent=2))  # Add indentation for readability (optional)

  except IOError as e:
    raise ValueError(f"Error opening or writing file: {e}")




##Using Modify_dict_file, setting args
filename = "./params/na16HH_TF2.txt"
changes = {
          "sh":8,
          "gbar":0.1,
          "tha":-59, #don't change
          "qa":4.5, #don't change
          "Ra":0.4,
          "Rb":0.124,
          "thi1":-80, #don't change
          "thi2":-80, #don't change
          "qd":5.4, #don't change
          "qg":5.4, #don't change
          "hmin":0.01,
          "mmin":0.02,
          "q10":2,
          "Rg":0.01,
          "Rd":0.03,
          "thinf":-80, #don't change
          "qinf":5.4, #don't change
          "vhalfs":-60, #don't change
          "a0s":0.0003,
          "zetas":12,
          "gms":0.2,
          "smax":10,
          "vvh":-58,
          "vvs":2,
          "ar2":1,
          "ena":55
          }

# modify_dict_file(filename, changes)

def combine_dictionaries(folder_path, new_file):
  """
  Combines dictionaries from .txt files in a folder into a single dictionary.

  Args:
      folder_path (str): Path to the folder containing .txt files with dictionaries.

  Returns:
      dict: A dictionary where file names are keys and dictionary contents are values.
  """
  combined_dict = {}
  for filename in os.listdir(folder_path):
    # Check for .txt files
    if filename.endswith(".txt"):
      file_path = os.path.join(folder_path, filename)
      try:
        # Open file and read content
        with open(file_path, 'r') as f:
          data = json.load(f)
        # Add dictionary to combined dict with filename as key
        combined_dict[filename] = data
      except (FileNotFoundError, json.JSONDecodeError) as e:
        print(f"Error processing file {filename}: {e}")
  with open (new_file ,'w') as file:
    # json.dump(combined_dict, file, separators=(',', ':')) #indent=4 if you want more vertical orientation
    # json.dump(combined_dict, file, indent=4) #indent=4 if you want more vertical orientation
    for key, value in combined_dict.items():
    # Write key-value pair on a separate line
      file.write(f'"{key}": {json.dumps(value)},\n')
  return combined_dict


## Plot states for 8 state HMM model
def plot_8states(csv_name,outfile_sfx,start=6500,stop=8500, ap_t=None, vm_t=None):
  df=pd.read_csv(csv_name, index_col=False)
  # Ensure start and end indices are within valid range
  if start < 0 or start >= len(df):
    raise ValueError("Invalid start index. Must be non-negative and less than data length.")
  if stop <= start or stop > len(df):
    raise ValueError("Invalid end index. Must be greater than start index and within data length.")
  
  dfsub = df.iloc[start:stop]

  fig1, ax1 = plt.subplots()
  ax1.plot(dfsub['c1'],label='c1', color='red', linewidth=0.7)
  ax1.plot(dfsub['c2'],label='c2', color='orange', linewidth=0.7)
  ax1.plot(dfsub['c3'],label='c3', color='pink', linewidth=0.7)
  ax1.plot(dfsub['i1'],label='i1', color='green', linewidth=0.7)
  ax1.plot(dfsub['i2'],label='i2', color='blue', linewidth=0.7)
  ax1.plot(dfsub['i3'],label='i3', color='cyan', linewidth=0.7)
  ax1.plot(dfsub['i4'],label='i4', color='purple', linewidth=0.7)
  ax1.plot(dfsub['o'],label='o', color='black', linewidth=0.7)

  ax2 = ax1.twinx()

  if ap_t is not None and vm_t is not None:
    y_min = min(np.min(vm_t[start:stop]), np.min(dfsub.min(axis=1)))
    y_max = max(np.max(vm_t[start:stop]), np.max(dfsub.max(axis=1)))
    ax2.set_ylim(bottom=y_min, top=y_max)
    ax2.plot(ap_t[start:stop],vm_t[start:stop], label='Vm', color='black',linewidth=2,linestyle='dashed') ##TF031424 changed linewidth
    ax2.set_ylabel('Vm', color='black')
  
  ax1.set_xlabel('timestep')
  ax1.set_ylabel('stateval')
  ax1.legend(loc='upper left')
  
  if ap_t is not None and vm_t is not None:
    ax2.legend(loc='upper right')
  
  # plt.legend()
  # plt.xticks(np.arange(0,len(df),1000), rotation=270)
  plt.xticks(np.arange(start,stop,500), rotation=270)
  plt.xlabel('timestep')
  # plt.ylabel('stateval')
  plt.title("States")
  # plt.savefig("8States"+".png", dpi=400)
  plt.savefig(f"/global/homes/t/tfenton/Neuron_general-2/Plots/Channel_state_plots/{start}-{stop}_{outfile_sfx}.png", dpi=400)


# combined_dict = combine_dictionaries(folder_path='/global/homes/t/tfenton/Neuron_general-2/params/na16_HOF_params_JSON', new_file='/global/homes/t/tfenton/Neuron_general-2/params/na16_HOF_params_JSON/combined3.json')

<<<<<<< HEAD
combined_dict = combine_dictionaries(folder_path='./params/na16_HOF_params_JSON', new_file='./params/na16_HOF_params_JSON/combined3.json')
=======
>>>>>>> d379827d8cc96912385ee0fb7f15fc15004d59a6

# plot_8states(csv_name="pandas_states.csv", outfile_sfx="8st_062824")



#combine_pdfs(folder_path='/global/homes/t/tfenton/Neuron_general-2/Plots/12HMM16HH_TF/AllSynthMuts_121223/', out_sfx='traceDVDTfi_121323')


#combine_and_sort_ef_csvs(root_path='/global/homes/t/tfenton/Neuron_general-2/Plots/12HMM16HH_TF/AllSynthMuts_121223/', out_sfx='allsynthmuts_121323_sorted')

# make_ppt_from_pdf2(pdf_path='/global/homes/t/tfenton/Neuron_general-2/Plots/12HMM16HH_TF/ManuscriptFigs/finetune_nav16',
                  # output_ppt_path='/global/homes/t/tfenton/Neuron_general-2/Plots/12HMM16HH_TF/ManuscriptFigs/finetune_nav16/finetune_na16_012524.pptx')

#plot_efeatures_bar(plot_folder='/global/homes/t/tfenton/Neuron_general-2/Plots/12HMM16HH_TF/ManuscriptFigs/efeatures',pfx='soma')




